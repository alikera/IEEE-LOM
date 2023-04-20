import json
import os
import zipfile
import imghdr
import shutil
import collections 
import collections.abc
from pptx import Presentation
import pptx
import filetype
import statistics
import spacy

import nltk
from nltk.corpus import wordnet
from nltk.stem import WordNetLemmatizer
import string
import textstat
# https://pypi.org/project/textstat/
# https://py-readability-metrics.readthedocs.io/en/latest/

from readability import Readability


import spacy
from spacy.language import Language
from spacy_langdetect import LanguageDetector


def insert_slide_initial_metadata(lom, resource_type, interactivity):
    lom['learningResourceType'] = resource_type
    lom['interactivityType'] = interactivity


def count_voices_in_slide(power):
    print('Extracting...')
    audio_extensions = ('.aiff', '.au', '.mid', '.midi', '.mp3', '.m4a', '.mp4', '.wav', '.wma')
    init_dir = os.getcwd()

    if not power.endswith('.zip') and not power.endswith('.pptx') or power == 'base_library.zip':
        print("here")
        return

    base = os.path.splitext(power)[0]
    new_zip = 'assets/temp/' + base + '.zip'
    os.rename('assets/temp/'+power, new_zip)
    with zipfile.ZipFile(new_zip) as myzip:
        if myzip.testzip() is not None:
            print("Some of your media files are corrupted")
        else:
            for file in myzip.namelist():
                if file.endswith(audio_extensions):
                    myzip.extract(file)

    myzip.close()

    try:
        os.chdir(init_dir + '/ppt/media')
        count = (len([name for name in os.listdir('.') if os.path.isfile(name)]))
        os.chdir(init_dir)
        shutil.rmtree(init_dir + '/ppt')
        return count
    except FileNotFoundError:
        return 0

def count_interactions_in_slide(power):
    ppt = Presentation("assets/"+power)
    print(power)
    num_videos = 0
    num_forms = 0
    num_clicks = 0
    num_hyperlinks = 0
    num_voices = 0
    counter = -1
    for slide in ppt.slides:
        counter += 1
        # print("slide:                  ",counter)
        for shape in slide.shapes:
            # if(type(shape) is pptx.shapes.autoshape.Shape):
                # print(type(shape), shape.shape_type)
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.MEDIA:
                num_videos += 1
                
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.FORM_CONTROL:
                num_forms += 1
            
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE:
                if(shape.click_action.target_slide is not None):
                    num_clicks += 1
                # print(shape.click_action._hlink)
                if(shape.click_action.hyperlink is not None):
                    num_hyperlinks += 1

    num_voices = count_voices_in_slide(power)

    points = 0
    if(num_videos > 0):
        points += 1
    if(num_voices > len(ppt.slides) /20):
        points += 1
    if(num_forms > 0):
        points += 1
    if(num_clicks > len(ppt.slides)/2):
        points += 1
    if(num_hyperlinks > len(ppt.slides) /10):
        points += 1
    


    # print("Number of videos:", num_videos)
    # print("Number of voices:", num_voices)

    # print("Number of forms:", num_forms)
    # print("Number of clicks:", num_clicks)
    # print("Number of hyperlinks:", num_hyperlinks)

    return points


def insert_interactivity_level(lom,count):
    lom["interactivityLevel"] = count


def process_slide_interactivity_level_tag(lom, slide):
    count = count_interactions_in_slide(slide)
    insert_interactivity_level(lom, count + 1)


def calculate_unique_words_density(text):
    # Initialize lemmatizer
    lemmatizer = WordNetLemmatizer()

    text_without_punctuations = text.translate(str.maketrans("", "", string.punctuation))

    # Tokenize the text into words
    words = nltk.word_tokenize(text_without_punctuations)

    # Remove stopwords
    stopwords = nltk.corpus.stopwords.words('english')
    words = [word for word in words if word.lower() not in stopwords]

    # Lemmatize the words
    lemmatized_words = []
    for word in words:
        # Get the part of speech (pos) tag for the word
        pos_tag = nltk.pos_tag([word])[0][1][0].lower()
        
        # Map the pos tag to the WordNet pos tags
        if pos_tag == 'j':
            pos_tag = wordnet.ADJ
        elif pos_tag == 'v':
            pos_tag = wordnet.VERB
        elif pos_tag == 'n':
            pos_tag = wordnet.NOUN
        elif pos_tag == 'r':
            pos_tag = wordnet.ADV
        else:
            pos_tag = None
        
        if pos_tag:
            # Lemmatize the word using the WordNet pos tag
            lemma = lemmatizer.lemmatize(word, pos=pos_tag)
        else:
            lemma = lemmatizer.lemmatize(word)
            
        lemmatized_words.append(lemma)

    # Remove duplicates
    unique_words = list(set(lemmatized_words))
    num_unique_words = len(unique_words)
    num_words = len(lemmatized_words)
    if num_words == 0:
        return None
    semantic_density = num_unique_words / num_words
    # Print the unique words after lemmatization and removing duplicates
    return semantic_density
    # print(semantic_density)


def calculate_slide_semantic_density(curr_slide):
    slide_text = ""
    each_slide_text = ""
    densities = []
    ppt = Presentation("assets/"+curr_slide)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += " " + run.text
                    each_slide_text += " " + run.text
        
        curr_density = calculate_unique_words_density(each_slide_text)
        if curr_density is not None:
            densities.append(curr_density)
        each_slide_text = ""
    
    each_slide_density = (statistics.median(densities) +statistics.mean(densities))/2
    semantic_density = (calculate_unique_words_density(slide_text) + each_slide_density )/2

    if(semantic_density <= 0.5):
        return 1
    if(semantic_density > 0.5 and semantic_density <= 0.6):
        return 2
    if(semantic_density > 0.6 and semantic_density <= 0.7):
        return 3
    if(semantic_density > 0.7 and semantic_density <= 0.8):
        return 4
    if(semantic_density > 0.8):
        return 5

    # print(calculate_unique_words_density(slide_text))
    # print(statistics.median(densities))
    # print(statistics.mean(densities))

def calculate_slide_semantic_density2(curr_slide):
    nlp = spacy.load('en_core_web_sm')
    slide_text = ""
    each_slide_text = ""
    ppt = Presentation("assets/"+curr_slide)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += " " + run.text
                    each_slide_text += " " + run.text

    doc = nlp(slide_text)
    similarity_scores = []
    for token1 in doc:
        similarity = 0
        for token2 in doc:
            if token1 != token2:
                similarity += token1.similarity(token2)
        similarity_scores.append(similarity)
    
    semantic_density = sum(similarity_scores) / len(similarity_scores)
    print(semantic_density)


def process_slide_semantic_density_tag(lom, slide):
    semantic_density_score = calculate_slide_semantic_density(slide)
    lom["semanticDensity"] = semantic_density_score


def process_slide_intended_end_user_role_tag(lom, slide):
    lom['intendedEndUserRole'] = 'learner'

def process_slide_context_tag(lom, slide):
    age = calculate_age_range(slide)
    if(age+1 >= 18):
        lom['context'] = 'higher education'
    elif(age+1 < 18):
        lom['context'] = 'school'
    

def extract_text_from_slide(ppt):
    text = ""
    # Iterate over each slide
    for slide in ppt.slides:
        # Iterate over each shape in the slide
        for shape in slide.shapes:
            # Check if the shape contains text
            if shape.has_text_frame:
                # Extract the text from the shape
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += "\n" + run.text
    
    return text

def calculate_slide_difficulty(curr_slide):
    # Open the PowerPoint presentation
    ppt = pptx.Presentation("assets/"+curr_slide)
    text = extract_text_from_slide(ppt)
    linsear_score = textstat.linsear_write_formula(text)
    # Print the text and the readability scores
    if(linsear_score < 5):
        return 1
    if(linsear_score >= 5 and linsear_score < 10):
        return 2
    if(linsear_score >= 10 and linsear_score < 15):
        return 3
    if(linsear_score >= 15 and linsear_score < 20):
        return 4
    if(linsear_score >= 20):
        return 5


def process_slide_difficulty_tag(lom, slide):
    difficulty = calculate_slide_difficulty(slide)
    lom['difficulty'] = difficulty
    

def calculate_slide_age_estimation(curr_slide):
    ppt = pptx.Presentation("assets/"+curr_slide)
    text = extract_text_from_slide(ppt)

    return (textstat.text_standard(text, float_output=True) + textstat.linsear_write_formula(text))/2


def calculate_age_range(slide):
    age_estimated = calculate_slide_age_estimation(slide)
    difficulty = calculate_slide_difficulty(slide)
    age_mid = int((age_estimated + difficulty)/2)

    return age_mid + 9
def process_slide_typical_age_range_tag(lom, slide):
    age = calculate_age_range(slide)
    lom['typicalAgeRange'] = str(age-1) + "-" + str(age+1)


def convert_seconds(seconds):
    hours, remainder = divmod(seconds, 3600)
    minutes, seconds = divmod(remainder, 60)
    return hours, minutes, seconds


def calculate_slide_learning_time(lom, curr_slide):
    ppt = pptx.Presentation("assets/"+curr_slide)
    text = extract_text_from_slide(ppt)
    new_dir = "extracted_context"
    filename = os.path.splitext(curr_slide)[0] + '.txt'
    filedir = os.path.join(new_dir, filename)
    with open(filedir, "w") as f:
        f.write(text)
    
    char_read_time = (((300 / lom['difficulty'])) / 60) * 5
    reading_time = textstat.reading_time(text, ms_per_char=1000/char_read_time)
    print(1000/char_read_time,reading_time)
    hours, minutes, seconds = convert_seconds(reading_time)

    return "PT" + str(int(hours)) + "H" + str(int(minutes)) + "M" + str(int(seconds)) + "S"


def process_slide_typical_learning_time(lom, slide):
    reading_time = calculate_slide_learning_time(lom, slide)
    print(reading_time)

    lom['typicalLearningTime'] = reading_time


def process_slide_description_tag(lom, slide):
    lom['description'] = ""

def get_lang_detector(nlp, name):
    return LanguageDetector()

def detect_language_text(curr_slide):
    ppt = pptx.Presentation("assets/"+curr_slide)
    text = extract_text_from_slide(ppt)

    nlp = spacy.load("en_core_web_sm")
    Language.factory("language_detector", func=get_lang_detector)
    nlp.add_pipe('language_detector', last=True)
    doc = nlp(text)

    detect_language = doc._.language #4

    return detect_language['language']

def process_slide_language_tag(lom, slide):
    language = detect_language_text(slide)
    print(language)
    lom['language'] = language


def copy_assets():
    init_dir = os.getcwd()

    origin = init_dir + '/assets/'
    target = origin + '/temp/'

    # Fetching the list of all the files
    files = os.listdir(origin)

    # Fetching all the files to directory
    for file_name in files:
        if file_name == 'temp':
            continue
        shutil.copy(origin+file_name, target+file_name)


